from pptx import Presentation
from pptx.util import Pt

SLD_LAYOUT_TITLE_AND_CONTENT = 1
prs = Presentation()
i=0
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
file= open("summary.txt","r")
content=file.readlines()
for a in content:
    text1= str(content)
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Presentation'
tf = body_shape.text_frame
mylist=[]
mylist = text1.split('. ') # splitting the entire data into single sentences and adding them to array..
count=1
for sentences in mylist: # iterate over each sentence to create a new bullet..
    p = tf.add_paragraph() # add paragraph is to add new bullet..
    run = p.add_run() # this code is to create a new run object related to that specific bullet
    run.text = sentences # adding text to run and not the paragraph..
    font = run.font # run lets us format the text...
    #font.name = 'Times New Roman'  or put 'Calibri' here, make sure spelling is correct! :D
    font.size = Pt(24) # set font size here...
    run.text = sentences
    if count==4: # creating a new slide after 4 bullets are on a slide, you can change the count if you want
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = 'Presentation'
        tf = body_shape.text_frame
        count=1 # resetting the count for the new slide
    count+=1
prs.save('testppt.pptx')
